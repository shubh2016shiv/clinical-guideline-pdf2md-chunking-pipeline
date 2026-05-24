"""
stage1_document_prescanning/document_page_structure_scanner.py
===============================================================
Module 2 of Stage 1: walk every page and build a structural fingerprint.

Responsibilities
----------------
For each page in the source document, produce a ``PageProfile`` — five numbers
that the complexity classifier uses to decide which extraction engine to invoke.

    page_number       int     1-based page index
    is_multi_column   bool    two or more text column bands detected
    has_diagrams      bool    many images with sparse text (likely a flowchart)
    has_large_tables  bool    text spans more than 60 % of the page width
    text_density      float   characters per mm² of page area (0.0 – 1.0 capped)

No GPU, no ML model, no external API is called here.  Everything is structural
arithmetic over the document's layout metadata.

Format coverage
---------------
PDF   → pypdfium2  (lightweight C binding, no rendering required)
DOCX  → python-docx (section columns, tables, paragraph images)
PPTX  → python-pptx (each slide is one page; shapes expose layout)
HTML  → stdlib html.parser (estimates from DOM element counts)

Text-density unit
-----------------
All formats normalise text_density to characters per mm² so the classifier
thresholds in settings.yaml are consistent across file types.

    text_density = char_count / page_area_mm²

Conversion factors:
  PDF/DOCX:  1 pt  = (25.4 / 72) mm  →  1 pt²  = (25.4/72)²  ≈ 0.1245 mm²
  PPTX:      1 EMU = (25.4/914400) mm →  1 EMU² = (25.4/914400)² mm²
  HTML:      uses standard A4 area (210 × 297 = 62 370 mm²) per estimated page

Approximate ranges for well-formed documents:
  Full-text page:   0.03 – 0.15 chars/mm²
  Sparse/mixed:     0.005 – 0.03 chars/mm²
  Diagram/blank:    0.0   – 0.005 chars/mm²

These ranges inform the has_diagrams and low_text_density thresholds used by
the complexity classifier (see document_complexity_classifier.py).
"""

from __future__ import annotations

import math
from dataclasses import dataclass
from html.parser import HTMLParser
from pathlib import Path
from typing import TYPE_CHECKING

import pypdfium2 as pdfium
import pypdfium2.raw as _pdfium_c

from ..contracts.configurations.pipeline_config import DocumentConstraintsConfig
from ..contracts.exceptions import DocumentError, DocumentTooLargeError
from ..contracts.pipeline_domain_types import DocumentType, PageProfile

if TYPE_CHECKING:
    pass

# ---------------------------------------------------------------------------
# Unit-conversion constants
# ---------------------------------------------------------------------------

# 1 PDF/typographic point = 25.4/72 mm ≈ 0.35278 mm
# 1 pt² ≈ 0.12445 mm²
_PT2_TO_MM2: float = (25.4 / 72.0) ** 2

# 1 EMU (English Metric Unit, PPTX) = 1/914400 inch = 25.4/914400 mm
# 1 EMU² ≈ 7.716e-10 mm²
_EMU_TO_MM: float = 25.4 / 914400.0
_EMU2_TO_MM2: float = _EMU_TO_MM**2

# Standard A4 page area used as a reference for HTML and DOCX pages
# whose actual dimensions cannot be read from the file itself.
_A4_AREA_MM2: float = 210.0 * 297.0  # = 62 370.0 mm²

# Minimum gap in the X-projection histogram that signals a column boundary.
# A gap of 10 % of page width between adjacent text-rectangle centres
# indicates two distinct text columns.  (Documented algorithm.)
_MULTI_COLUMN_GAP_RATIO: float = 0.10

# A text rectangle whose width exceeds this fraction of the page width is
# classified as a potential large table spanning the full content area.
_LARGE_TABLE_WIDTH_RATIO: float = 0.60

# Images per page threshold used together with text_density in has_diagrams.
_DIAGRAM_IMAGE_COUNT_THRESHOLD: int = 2

# text_density threshold (chars/mm²) below which a page is diagram-heavy.
# Calibrated so that a page with only a few caption lines (< ~3 chars/mm²
# across a 150 mm² area) triggers the flag when combined with image count.
# Adjust in the scanner if your corpus skews toward denser or sparser pages.
_DIAGRAM_DENSITY_THRESHOLD: float = 0.05

# Estimated characters per page for DOCX/HTML (used when actual rendering
# info is unavailable).  300 words × 5.5 chars/word ≈ 1 650; round to 2 000
# to be conservative for clinical guidelines with tables and whitespace.
_ESTIMATED_CHARS_PER_PAGE: int = 2_000

# For HTML page estimation: 1 page ≈ _ESTIMATED_CHARS_PER_PAGE characters.
# For DOCX page estimation: same.


@dataclass(frozen=True)
class DocumentStructureScanResult:
    """
    Output produced by ``DocumentPageStructureScanner.scan()``.

    ``total_pages`` is set on ``ConversionJob`` after Stage 1 completes.
    ``profiles`` is the sole input to the complexity classifier.
    """

    total_pages: int
    """Actual page count (PDF, PPTX) or estimated page count (DOCX, HTML)."""

    profiles: list[PageProfile]
    """One PageProfile per page, in document order."""


class DocumentPageStructureScanner:
    """
    Walks every page of a document and emits a ``PageProfile`` for each.

    The scanner is format-aware and dispatches to the lightest possible
    library for each format — no rendering, no GPU, no ML.

    Instantiate once with the configured constraints::

        scanner = DocumentPageStructureScanner(config.document_constraints)
        result = scanner.scan(path, DocumentType.PDF)
        # result.total_pages → 100
        # result.profiles    → [PageProfile(...), ...]
    """

    def __init__(self, constraints: DocumentConstraintsConfig) -> None:
        self._max_pages = constraints.max_pages

    def scan(self, document_path: Path, document_type: DocumentType) -> DocumentStructureScanResult:
        """
        Scan all pages of *document_path* and return structural profiles.

        Args:
            document_path: Absolute path to the source document (already validated
                           by ``DocumentSHA256Hasher``).
            document_type: Format of the document (from the hasher).

        Returns:
            ``DocumentStructureScanResult`` with total page count and profiles.

        Raises:
            DocumentError: Document is corrupt, has zero pages, or the format
                is not supported by this scanner.
            DocumentTooLargeError: Page count exceeds ``max_pages``.
        """
        dispatch = {
            DocumentType.PDF: self._scan_pdf,
            DocumentType.DOCX: self._scan_docx,
            DocumentType.PPTX: self._scan_pptx,
            DocumentType.HTML: self._scan_html,
        }
        scan_fn = dispatch.get(document_type)
        if scan_fn is None:
            raise DocumentError(
                f"Stage 1 structure scanning is not supported for {document_type.value!r}.",
                context={"document_type": document_type.value, "path": str(document_path)},
            )
        return scan_fn(document_path)

    # ------------------------------------------------------------------
    # PDF scanner (pypdfium2)
    # ------------------------------------------------------------------

    def _scan_pdf(self, path: Path) -> DocumentStructureScanResult:
        """
        Open the PDF with pypdfium2 and walk every page.

        pypdfium2 reads text-block metadata directly from the PDF's internal
        structure — it does NOT render pages to bitmaps.  That is why 500 pages
        takes under 2 seconds even on a modest CPU.
        """
        try:
            doc = pdfium.PdfDocument(str(path))
        except Exception as exc:
            raise DocumentError(
                f"Failed to open PDF (corrupt or invalid): {path.name}",
                context={"path": str(path)},
            ) from exc

        total_pages = len(doc)
        if total_pages == 0:
            raise DocumentError(
                f"PDF has zero pages: {path.name}",
                context={"path": str(path)},
            )
        self._assert_page_limit(total_pages, path)

        profiles: list[PageProfile] = []
        for page_idx in range(total_pages):
            page = doc[page_idx]
            profiles.append(self._profile_pdf_page(page, page_idx + 1))

        return DocumentStructureScanResult(total_pages=total_pages, profiles=profiles)

    def _profile_pdf_page(self, page: pdfium.PdfPage, page_number: int) -> PageProfile:
        """Compute one PageProfile from a pypdfium2 PdfPage."""
        page_width, page_height = page.get_size()

        # Guard against zero-area pages (malformed PDFs occasionally have them).
        if page_width <= 0 or page_height <= 0:
            return PageProfile(
                page_number=page_number,
                is_multi_column=False,
                has_diagrams=False,
                has_large_tables=False,
                text_density=0.0,
            )

        # --- text extraction -----------------------------------------------
        textpage = page.get_textpage()

        # Full-page text for density calculation.
        full_text = textpage.get_text_bounded()
        char_count = len(full_text)

        # Text rectangles — PDFium groups nearby characters into rectangular
        # runs.  These are our "text blocks" for layout analysis.
        # count_rects() MUST be called once before any get_rect() calls
        # (it initialises PDFium's internal iterator state).
        n_rects = textpage.count_rects()
        rects: list[tuple[float, float, float, float]] = [
            textpage.get_rect(i) for i in range(n_rects)
        ]  # each rect is (left, bottom, right, top)

        # --- image counting ------------------------------------------------
        # filter= accepts a list of FPDF_PAGEOBJ_* integer constants.
        image_objects = list(page.get_objects(filter=[_pdfium_c.FPDF_PAGEOBJ_IMAGE]))
        image_count = len(image_objects)

        # --- derived flags -------------------------------------------------
        page_area_mm2 = page_width * page_height * _PT2_TO_MM2
        text_density = min(char_count / max(page_area_mm2, 1.0), 1.0)

        is_multi_column = _detect_multi_column(rects, page_width)
        has_large_tables = _detect_large_tables(rects, page_width)
        has_diagrams = (
            image_count > _DIAGRAM_IMAGE_COUNT_THRESHOLD
            and text_density < _DIAGRAM_DENSITY_THRESHOLD
        )

        return PageProfile(
            page_number=page_number,
            is_multi_column=is_multi_column,
            has_diagrams=has_diagrams,
            has_large_tables=has_large_tables,
            text_density=text_density,
        )

    # ------------------------------------------------------------------
    # DOCX scanner (python-docx)
    # ------------------------------------------------------------------

    def _scan_docx(self, path: Path) -> DocumentStructureScanResult:
        """
        Analyse a DOCX document with python-docx.

        DOCX has no native page-count concept — pages are determined by a
        layout engine (Word, LibreOffice) that we do not invoke here.  Instead
        we estimate page count from total character count (÷ _ESTIMATED_CHARS_PER_PAGE).

        For complexity profiling this is sufficient: a 100-page clinical
        guideline and a 5-page summary will have very different raw character
        counts, and the classifier only needs the proportion of complex pages.
        """
        try:
            from docx import Document as DocxDocument  # python-docx (transitive via docling)
        except ImportError as exc:
            raise DocumentError(
                "python-docx is required to pre-scan DOCX documents.  "
                "Install it with: pip install python-docx",
                context={"path": str(path)},
            ) from exc

        try:
            doc = DocxDocument(str(path))
        except Exception as exc:
            raise DocumentError(
                f"Failed to open DOCX (corrupt or password-protected): {path.name}",
                context={"path": str(path)},
            ) from exc

        # --- document-level features --------------------------------------

        # Column layout: check if any section uses more than one column.
        is_multi_column = _docx_has_multi_column_section(doc)

        # Page dimensions from the first section (used for text_density).
        first_section = doc.sections[0] if doc.sections else None
        if first_section is not None:
            try:
                page_area_mm2 = (
                    first_section.page_width.pt * first_section.page_height.pt * _PT2_TO_MM2
                )
            except Exception:
                page_area_mm2 = _A4_AREA_MM2
        else:
            page_area_mm2 = _A4_AREA_MM2

        # Aggregate character count across all paragraphs and table cells.
        total_chars = sum(len(p.text) for p in doc.paragraphs)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    total_chars += sum(len(p.text) for p in cell.paragraphs)

        table_count = len(doc.tables)
        image_count = _count_docx_inline_images(doc)

        # --- page estimation -----------------------------------------------
        estimated_pages = max(1, math.ceil(total_chars / _ESTIMATED_CHARS_PER_PAGE))
        self._assert_page_limit(estimated_pages, path)

        chars_per_page = total_chars / estimated_pages
        text_density = min(chars_per_page / max(page_area_mm2, 1.0), 1.0)
        images_per_page = image_count / estimated_pages

        # Build one profile per estimated page.  All pages receive the same
        # document-level features because DOCX does not encode per-page layout.
        has_large_tables = table_count > 0  # any table may span full page width
        has_diagrams = (
            images_per_page > _DIAGRAM_IMAGE_COUNT_THRESHOLD
            and text_density < _DIAGRAM_DENSITY_THRESHOLD
        )

        profiles = [
            PageProfile(
                page_number=i + 1,
                is_multi_column=is_multi_column,
                has_diagrams=has_diagrams,
                has_large_tables=has_large_tables,
                text_density=text_density,
            )
            for i in range(estimated_pages)
        ]

        return DocumentStructureScanResult(total_pages=estimated_pages, profiles=profiles)

    # ------------------------------------------------------------------
    # PPTX scanner (python-pptx)
    # ------------------------------------------------------------------

    def _scan_pptx(self, path: Path) -> DocumentStructureScanResult:
        """
        Analyse a PPTX presentation with python-pptx.

        Each slide is treated as one page.  Shape bounding boxes are in EMU
        (English Metric Units); they are converted to mm for density calculations.
        """
        try:
            from pptx import Presentation  # python-pptx (transitive via docling)
            from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: PLC0415
        except ImportError as exc:
            raise DocumentError(
                "python-pptx is required to pre-scan PPTX documents.  "
                "Install it with: pip install python-pptx",
                context={"path": str(path)},
            ) from exc

        try:
            prs = Presentation(str(path))
        except Exception as exc:
            raise DocumentError(
                f"Failed to open PPTX (corrupt or password-protected): {path.name}",
                context={"path": str(path)},
            ) from exc

        slides = list(prs.slides)
        total_slides = len(slides)
        if total_slides == 0:
            raise DocumentError(
                f"PPTX has zero slides: {path.name}",
                context={"path": str(path)},
            )
        self._assert_page_limit(total_slides, path)

        slide_width_emu: int = prs.slide_width or 9_144_000  # default 10 inches
        slide_height_emu: int = prs.slide_height or 5_143_500  # default 7.5 inches
        slide_area_mm2 = slide_width_emu * slide_height_emu * _EMU2_TO_MM2

        profiles: list[PageProfile] = []
        for slide_idx, slide in enumerate(slides):
            shapes = list(slide.shapes)

            image_count = sum(1 for s in shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
            char_count = sum(len(s.text_frame.text) for s in shapes if s.has_text_frame)
            # Large-table heuristic: any table shape wider than 60 % of the slide.
            # Filter to table shapes only, and guard against None width values
            # (PPTX shapes can report width=None if not explicitly set).
            table_shapes = [s for s in shapes if s.has_table]
            has_large_tables = any(
                s.width is not None
                and (s.width / max(slide_width_emu, 1)) > _LARGE_TABLE_WIDTH_RATIO
                for s in table_shapes
            )

            text_density = min(char_count / max(slide_area_mm2, 1.0), 1.0)
            has_diagrams = (
                image_count > _DIAGRAM_IMAGE_COUNT_THRESHOLD
                and text_density < _DIAGRAM_DENSITY_THRESHOLD
            )

            # Multi-column detection: use the same X-projection algorithm as PDF
            # but on text frame (placeholder + content box) left edges in EMU.
            text_box_rects = [
                (s.left, 0, s.left + (s.width or 0), 1)
                for s in shapes
                if s.has_text_frame and not s.has_table and s.left is not None
            ]
            is_multi_column = _detect_multi_column(text_box_rects, float(slide_width_emu))

            profiles.append(
                PageProfile(
                    page_number=slide_idx + 1,
                    is_multi_column=is_multi_column,
                    has_diagrams=has_diagrams,
                    has_large_tables=has_large_tables,
                    text_density=text_density,
                )
            )

        return DocumentStructureScanResult(total_pages=total_slides, profiles=profiles)

    # ------------------------------------------------------------------
    # HTML scanner (stdlib html.parser)
    # ------------------------------------------------------------------

    def _scan_html(self, path: Path) -> DocumentStructureScanResult:
        """
        Estimate structural complexity of an HTML document.

        HTML has no inherent page concept.  We parse the DOM to count text
        characters, images, and tables, then divide by _ESTIMATED_CHARS_PER_PAGE
        to produce a page-count estimate.  One A4-equivalent page area is used
        for text_density normalisation.

        Clinical guidelines rarely arrive as raw HTML, but when they do (e.g.,
        exported from a CMS), the classifier needs something reasonable to route
        them to the right engine.
        """
        try:
            html_source = path.read_text(encoding="utf-8", errors="ignore")
        except Exception as exc:
            raise DocumentError(
                f"Failed to read HTML document: {path.name}",
                context={"path": str(path)},
            ) from exc

        extractor = _HTMLStructureExtractor()
        extractor.feed(html_source)

        total_chars = len(extractor.text)
        estimated_pages = max(1, math.ceil(total_chars / _ESTIMATED_CHARS_PER_PAGE))
        self._assert_page_limit(estimated_pages, path)

        chars_per_page = total_chars / estimated_pages
        text_density = min(chars_per_page / _A4_AREA_MM2, 1.0)

        images_per_page = extractor.image_count / estimated_pages
        has_diagrams = (
            images_per_page > _DIAGRAM_IMAGE_COUNT_THRESHOLD
            and text_density < _DIAGRAM_DENSITY_THRESHOLD
        )
        has_large_tables = extractor.table_count > 0
        # HTML column layout detection is impractical without CSS rendering;
        # default to False and let the engine handle it.
        is_multi_column = False

        profiles = [
            PageProfile(
                page_number=i + 1,
                is_multi_column=is_multi_column,
                has_diagrams=has_diagrams,
                has_large_tables=has_large_tables,
                text_density=text_density,
            )
            for i in range(estimated_pages)
        ]

        return DocumentStructureScanResult(total_pages=estimated_pages, profiles=profiles)

    # ------------------------------------------------------------------
    # Shared guard
    # ------------------------------------------------------------------

    def _assert_page_limit(self, page_count: int, path: Path) -> None:
        if page_count > self._max_pages:
            raise DocumentTooLargeError(
                f"Document has {page_count} pages which exceeds the configured "
                f"limit of {self._max_pages} pages",
                context={
                    "path": str(path),
                    "page_count": page_count,
                    "limit": self._max_pages,
                },
            )


# ---------------------------------------------------------------------------
# Layout-analysis helpers (format-agnostic)
# ---------------------------------------------------------------------------


def _detect_multi_column(
    rects: list[tuple[float, float, float, float]],
    container_width: float,
) -> bool:
    """
    X-projection gap detection for multi-column layout.

    Algorithm
    ---------
    1. Compute the horizontal centre of every text rectangle.
    2. Sort those centre values.
    3. Scan for a gap between any two consecutive sorted values that is at least
       ``_MULTI_COLUMN_GAP_RATIO`` × ``container_width`` wide.

    A gap of that size means there are two distinct bands of text separated by
    a gutter — the hallmark of a multi-column layout.  A single-column page
    has all text concentrated in one band with no wide gap.

    Time complexity: O(n log n) where n = number of text rectangles per page
    (typically 5–50 on a clinical guideline page, so effectively O(1)).

    Args:
        rects: List of (left, bottom, right, top) bounding boxes in any
               consistent unit (pts or EMU — the ratio is unit-agnostic).
        container_width: Page or slide width in the same unit as the rect coords.

    Returns:
        True when at least two distinct text columns are detected.
    """
    if not rects or container_width <= 0:
        return False

    # Centre X of every non-degenerate rectangle.
    centres = sorted((left + right) / 2.0 for left, _bottom, right, _top in rects if right > left)
    if len(centres) < 2:
        return False

    gap_threshold = container_width * _MULTI_COLUMN_GAP_RATIO

    # A single gap between any two consecutive sorted centres is enough
    # to confirm two distinct column bands.
    return any((centres[i + 1] - centres[i]) >= gap_threshold for i in range(len(centres) - 1))


def _detect_large_tables(
    rects: list[tuple[float, float, float, float]],
    page_width: float,
) -> bool:
    """
    Return True when any text rectangle spans more than 60 % of the page width.

    Wide text runs are the structural signature of full-width comparison tables
    common in clinical guidelines (e.g., "Drug A vs Drug B across 8 columns").
    A normal paragraph column uses ~40 % of a standard A4 page width.

    Args:
        rects: (left, bottom, right, top) bounding boxes in PDF canvas units.
        page_width: Total page width in the same unit as the rect coords.
    """
    if not rects or page_width <= 0:
        return False

    return any(
        (right - left) / page_width > _LARGE_TABLE_WIDTH_RATIO
        for left, _bottom, right, _top in rects
        if right > left
    )


# ---------------------------------------------------------------------------
# DOCX helpers
# ---------------------------------------------------------------------------


def _docx_has_multi_column_section(doc: object) -> bool:
    """
    Return True if any DOCX section is formatted with more than one text column.

    python-docx exposes the raw XML via ``section._sectPr``.  The ``<w:cols>``
    element holds the column count in the ``w:num`` attribute (default = 1).
    """
    try:
        from docx.oxml.ns import qn  # noqa: PLC0415
    except ImportError:
        return False

    for section in doc.sections:
        try:
            cols_el = section._sectPr.find(qn("w:cols"))
            if cols_el is not None:
                num = int(cols_el.get(qn("w:num"), "1"))
                if num > 1:
                    return True
        except Exception:
            # Malformed XML in a section should not abort the scan.
            continue
    return False


def _count_docx_inline_images(doc: object) -> int:
    """
    Count inline images embedded in paragraph runs of a DOCX document.

    python-docx does not expose a top-level ``doc.images`` collection; images
    are referenced from within run XML as ``<a:blip>`` elements inside picture
    shapes (``<pic:pic>``).  We search for the pic namespace element directly.
    """
    ns = "http://schemas.openxmlformats.org/drawingml/2006/picture"
    count = 0
    for para in doc.paragraphs:
        for run in para.runs:
            # Each run element may contain a drawing with one or more pictures.
            count += len(run._element.findall(f".//{{{ns}}}pic"))
    return count


# ---------------------------------------------------------------------------
# HTML structure extractor (stdlib only)
# ---------------------------------------------------------------------------


class _HTMLStructureExtractor(HTMLParser):
    """
    Minimal SAX-style parser that counts text content, images, and tables.

    Uses Python's built-in ``html.parser`` — no third-party dependency.
    Script and style content is excluded from the text count to avoid inflating
    character counts with CSS rules and JavaScript.
    """

    def __init__(self) -> None:
        super().__init__()
        self._skip_content: bool = False  # True while inside <script> or <style>
        self.text: list[str] = []
        self.image_count: int = 0
        self.table_count: int = 0

    def handle_starttag(self, tag: str, attrs: list) -> None:
        tag_lower = tag.lower()
        if tag_lower in ("script", "style"):
            self._skip_content = True
        elif tag_lower == "img":
            self.image_count += 1
        elif tag_lower == "table":
            self.table_count += 1

    def handle_endtag(self, tag: str) -> None:
        if tag.lower() in ("script", "style"):
            self._skip_content = False

    def handle_data(self, data: str) -> None:
        if not self._skip_content:
            stripped = data.strip()
            if stripped:
                self.text.append(stripped)

    @property
    def full_text(self) -> str:
        return " ".join(self.text)
