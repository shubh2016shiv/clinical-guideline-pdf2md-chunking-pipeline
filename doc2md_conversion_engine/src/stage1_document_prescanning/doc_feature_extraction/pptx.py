"""
PPTX feature extraction for engine routing.

This module does not try to understand the clinical meaning of a PowerPoint
presentation.  It only collects factual evidence that is cheap to inspect:

1. Is there native text?
2. Are there table shapes?
3. Are there native chart objects (bar, line, pie, etc.)?
4. Are there embedded raster images?
5. Are there drawing or connector shapes (lines, rectangles, arrows)?
6. Are there slides so visually dense they should be treated as diagrams?

Thresholds come from ``settings.yaml`` via ``DocumentFeatureExtractionConfig``.

-----

What is a PPTX file?
--------------------
Like DOCX, a .pptx file is a ZIP archive of XML files.  Each slide is a
separate XML file (slide1.xml, slide2.xml, ...) containing a flat list of
"shapes".  A shape is the PowerPoint term for any object placed on a slide —
a text box, an image, a chart, a table, a drawn line, a rectangle, etc.
Every shape has a type attribute that tells us what kind of object it is.

Why does slide area matter?
    Shapes store their size in EMUs (English Metric Units — PowerPoint's
    internal unit, where 914400 EMUs = 1 inch).  We convert shape size to a
    fraction of the total slide area so that comparisons are consistent
    regardless of whether the presentation uses widescreen (16:9) or
    standard (4:3) dimensions.
"""

from __future__ import annotations

import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from ...contracts.configurations.pipeline_config import (
    DocumentFeatureExtractionConfig,
    PptxFeatureExtractionConfig,
)
from ...contracts.exceptions import DocumentError
from .capabilities import get_engine_format_support
from .models import (
    DocumentFeatureProfile,
    FeatureDocumentType,
    TableEvidence,
    TextEvidence,
    VisualCandidate,
    VisualCandidateKind,
    VisualEvidence,
)
from .requirement_inference import infer_requirements
from .text_patterns import compact_text, contains_figure_caption, count_figure_caption_lines

logger = logging.getLogger(__name__)


@dataclass
class PptxFeatureTotals:
    """Running totals accumulated while scanning all slides in a presentation."""

    total_characters: int = 0
    slide_numbers_with_text: set[int] = field(default_factory=set)

    number_of_tables: int = 0
    slide_numbers_with_tables: set[int] = field(default_factory=set)

    number_of_images: int = 0
    number_of_large_images: int = 0
    number_of_charts: int = 0
    number_of_drawing_shapes: int = 0
    slide_numbers_with_visuals: set[int] = field(default_factory=set)

    number_of_captioned_figures: int = 0

    visual_routing_candidates: list[VisualCandidate] = field(default_factory=list)


@dataclass
class PptxSlideContext:
    """Mutable state accumulated while scanning the shapes on one slide."""

    slide_number: int
    text_blocks: list[str] = field(default_factory=list)
    visual_shape_count: int = 0


# ---------------------------------------------------------------------------
# Public entry-point
# ---------------------------------------------------------------------------


def extract_pptx_features(
    path: Path,
    config: DocumentFeatureExtractionConfig | None = None,
) -> DocumentFeatureProfile:
    """
    Extract deterministic feature evidence from one PowerPoint presentation.

    The flow is intentionally linear:
    open the file, visit each slide, collect evidence from each shape, then
    build the final ``DocumentFeatureProfile`` consumed by the capability router.

    Step 1: Open the file and read slide dimensions.
    Step 2: For each slide, scan every shape on it (text, tables, charts,
            images, drawings).
    Step 3: After all shapes on a slide, record per-slide text evidence and
            flag diagram-heavy slides.
    Step 4: Assemble and return the feature profile.
    """
    pptx_settings = (config or DocumentFeatureExtractionConfig()).pptx
    presentation, mso_shape_type = open_presentation(path)

    # Step 1: Slide dimensions — stored in EMUs; we only need the product
    # (slide area) to compute each shape's relative size later.
    slide_area_in_emus = max(
        float(presentation.slide_width or 0) * float(presentation.slide_height or 0),
        1.0,
    )
    total_slides = max(len(presentation.slides), 1)
    totals = PptxFeatureTotals()

    # Step 2 + 3: Walk every slide
    for slide_number, slide in enumerate(presentation.slides, start=1):
        scan_one_slide(
            slide=slide,
            slide_number=slide_number,
            slide_area_in_emus=slide_area_in_emus,
            settings=pptx_settings,
            mso_shape_type=mso_shape_type,
            totals=totals,
        )

    logger.debug(
        "PPTX feature extraction complete: path=%s slides=%d chars=%d "
        "tables=%d images=%d charts=%d drawings=%d",
        path.name,
        total_slides,
        totals.total_characters,
        totals.number_of_tables,
        totals.number_of_images,
        totals.number_of_charts,
        totals.number_of_drawing_shapes,
    )

    # Step 4: Assemble
    return build_pptx_feature_profile(
        total_slides=total_slides,
        totals=totals,
        settings=pptx_settings,
    )


# ---------------------------------------------------------------------------
# Step 2+3: Slide scanning — one slide at a time
# ---------------------------------------------------------------------------


def scan_one_slide(
    *,
    slide: Any,
    slide_number: int,
    slide_area_in_emus: float,
    settings: PptxFeatureExtractionConfig,
    mso_shape_type: Any,
    totals: PptxFeatureTotals,
) -> None:
    """
    Scan one slide in the same order a person would inspect it.

    Step 2a: Collect text from text-frame shapes (text boxes, titles, captions).
    Step 2b: Detect and record table shapes.
    Step 2c: Detect and record native chart objects.
    Step 2d: Detect and record picture (raster image) shapes.
    Step 2e: Detect and record drawing and connector shapes.
    Step 3:  After all shapes — record slide text evidence and flag diagram-heavy slides.
    """
    ctx = PptxSlideContext(slide_number=slide_number)

    for shape in slide.shapes:
        collect_text_from_shape(shape, ctx, totals)
        collect_table_from_shape(shape, slide_area_in_emus, ctx, totals)
        collect_chart_from_shape(shape, slide_area_in_emus, ctx, totals)
        collect_image_from_shape(shape, slide_area_in_emus, settings, mso_shape_type, ctx, totals)
        collect_drawing_from_shape(shape, mso_shape_type, ctx, totals)

    finalize_slide(ctx, settings, totals)


def collect_text_from_shape(
    shape: Any,
    ctx: PptxSlideContext,
    totals: PptxFeatureTotals,
) -> None:
    """
    Step 2a: Read text from any shape that carries a text frame.

    Text frames appear in text boxes, title placeholders, subtitle placeholders,
    bullet-point lists, and caption boxes — all the same API in python-pptx.
    """
    if not getattr(shape, "has_text_frame", False):
        return
    shape_text = getattr(shape, "text", "") or ""
    totals.total_characters += len(shape_text)
    ctx.text_blocks.append(shape_text)


def collect_table_from_shape(
    shape: Any,
    slide_area_in_emus: float,
    ctx: PptxSlideContext,
    totals: PptxFeatureTotals,
) -> None:
    """Step 2b: Record a table shape and emit a visual candidate for it."""
    if not getattr(shape, "has_table", False):
        return
    totals.number_of_tables += 1
    totals.slide_numbers_with_tables.add(ctx.slide_number)
    totals.visual_routing_candidates.append(
        VisualCandidate(
            kind=VisualCandidateKind.TABLE,
            page_number=ctx.slide_number,
            area_ratio=calculate_shape_area_as_fraction_of_slide(shape, slide_area_in_emus),
            nearby_text=find_caption_in_slide_text(ctx.text_blocks),
            evidence=["PowerPoint table shape"],
        )
    )


def collect_chart_from_shape(
    shape: Any,
    slide_area_in_emus: float,
    ctx: PptxSlideContext,
    totals: PptxFeatureTotals,
) -> None:
    """
    Step 2c: Record a native PowerPoint chart object.

    A "chart" in PPTX is a bar/line/pie chart linked to embedded spreadsheet
    data — distinct from a screenshot of a chart, which would be a picture shape.
    """
    if not getattr(shape, "has_chart", False):
        return
    totals.number_of_charts += 1
    totals.slide_numbers_with_visuals.add(ctx.slide_number)
    ctx.visual_shape_count += 1
    totals.visual_routing_candidates.append(
        VisualCandidate(
            kind=VisualCandidateKind.CHART,
            page_number=ctx.slide_number,
            area_ratio=calculate_shape_area_as_fraction_of_slide(shape, slide_area_in_emus),
            nearby_text=find_caption_in_slide_text(ctx.text_blocks),
            evidence=["PowerPoint native chart object"],
        )
    )


def collect_image_from_shape(
    shape: Any,
    slide_area_in_emus: float,
    settings: PptxFeatureExtractionConfig,
    mso_shape_type: Any,
    ctx: PptxSlideContext,
    totals: PptxFeatureTotals,
) -> None:
    """
    Step 2d: Record a raster picture shape.

    ``MSO_SHAPE_TYPE.PICTURE`` identifies JPEGs, PNGs, and similar embedded
    bitmaps inserted via Insert → Image.  We check whether the image is "large"
    (covers at least ``settings.large_visual_area_ratio`` of the slide) because
    large images are more likely to be primary content than decorative elements.
    """
    if shape.shape_type != mso_shape_type.PICTURE:
        return

    totals.number_of_images += 1
    totals.slide_numbers_with_visuals.add(ctx.slide_number)
    ctx.visual_shape_count += 1

    area_fraction = calculate_shape_area_as_fraction_of_slide(shape, slide_area_in_emus)
    if area_fraction is not None and area_fraction >= settings.large_visual_area_ratio:
        totals.number_of_large_images += 1

    totals.visual_routing_candidates.append(
        VisualCandidate(
            kind=VisualCandidateKind.EMBEDDED_IMAGE,
            page_number=ctx.slide_number,
            area_ratio=area_fraction,
            nearby_text=find_caption_in_slide_text(ctx.text_blocks),
            evidence=["PowerPoint picture shape (raster image)"],
        )
    )


def collect_drawing_from_shape(
    shape: Any,
    mso_shape_type: Any,
    ctx: PptxSlideContext,
    totals: PptxFeatureTotals,
) -> None:
    """
    Step 2e: Record a drawing or connector shape (line, rectangle, arrow, group).

    Drawing shapes are purely visual — they carry no structured content of their
    own.  We count them toward the slide's visual-shape total so that slides
    with many drawings can be flagged as diagram-heavy in Step 3.
    """
    if not shape_is_a_drawing_or_connector(shape, mso_shape_type):
        return
    totals.number_of_drawing_shapes += 1
    totals.slide_numbers_with_visuals.add(ctx.slide_number)
    ctx.visual_shape_count += 1


def finalize_slide(
    ctx: PptxSlideContext,
    settings: PptxFeatureExtractionConfig,
    totals: PptxFeatureTotals,
) -> None:
    """
    Step 3: Record per-slide summary evidence after all shapes are processed.

    Records whether this slide has any text and counts figure captions.
    Also flags the slide as a SLIDE_VISUAL_CLUSTER candidate when the number of
    visual shapes meets ``settings.diagram_heavy_slide_min_visual_shapes`` —
    that threshold comes from
    ``document_feature_extraction.pptx.diagram_heavy_slide_min_visual_shapes``
    in ``settings.yaml``.
    """
    full_slide_text = "\n".join(ctx.text_blocks)

    if full_slide_text.strip():
        totals.slide_numbers_with_text.add(ctx.slide_number)

    totals.number_of_captioned_figures += count_figure_caption_lines(full_slide_text)

    if ctx.visual_shape_count >= settings.diagram_heavy_slide_min_visual_shapes:
        totals.visual_routing_candidates.append(
            VisualCandidate(
                kind=VisualCandidateKind.SLIDE_VISUAL_CLUSTER,
                page_number=ctx.slide_number,
                nearby_text=find_caption_in_slide_text(ctx.text_blocks),
                evidence=[
                    f"{ctx.visual_shape_count} visual shapes on this slide; "
                    f"configured threshold is {settings.diagram_heavy_slide_min_visual_shapes}"
                ],
            )
        )


# ---------------------------------------------------------------------------
# Step 4: Profile assembly
# ---------------------------------------------------------------------------


def build_pptx_feature_profile(
    *,
    total_slides: int,
    totals: PptxFeatureTotals,
    settings: PptxFeatureExtractionConfig,
) -> DocumentFeatureProfile:
    """
    Translate raw extraction totals into the structured profile consumed by the
    capability router downstream.

    ``DocumentFeatureProfile`` is the single object that leaves this module.
    Everything else here is an intermediate view of the same raw numbers,
    shaped into the three distinct contracts the router expects:

    - *evidence* objects (TextEvidence, TableEvidence, VisualEvidence) — aggregate
      statistics (counts, per-slide flags).  The router reads these to decide
      *whether* a capability is needed.

    - *visual_candidates* — a ranked short-list of specific elements (an image
      shape, a chart, a table, or a diagram-heavy slide cluster) worth showing
      to a vision model for a second opinion.  Candidates carry the slide number,
      area fraction, and nearby caption text as context.  Capped by
      ``settings.max_visual_candidates`` so the downstream payload stays bounded.

    - *requirements* — inferred capability flags (e.g. needs_ocr, needs_vlm)
      derived from the evidence above.  The router uses these to filter the
      engine candidates it will score.
    """
    # Aggregate statistics — one object per feature dimension.
    text_evidence = build_text_evidence(totals, total_slides)
    table_evidence = build_table_evidence(totals)
    visual_evidence = build_visual_evidence(totals)

    # Representative examples — ranked by importance, capped so the downstream
    # router is not overwhelmed by a visually dense presentation.
    visual_candidates = choose_visual_candidates_for_routing(
        totals.visual_routing_candidates,
        max_candidates_to_keep=settings.max_visual_candidates,
    )

    # Capability flags — derived from evidence, not raw totals.
    # infer_requirements encodes the routing heuristics so this function stays
    # focused on assembly.
    requirements = infer_requirements(
        text=text_evidence,
        tables=table_evidence,
        visuals=visual_evidence,
        visual_candidates=visual_candidates,
    )

    return DocumentFeatureProfile(
        file_type=FeatureDocumentType.PPTX,
        page_or_unit_count=total_slides,
        text=text_evidence,
        tables=table_evidence,
        visuals=visual_evidence,
        visual_candidates=visual_candidates,
        format_support=get_engine_format_support(FeatureDocumentType.PPTX),
        requirements=requirements,
    )


def build_text_evidence(totals: PptxFeatureTotals, total_slides: int) -> TextEvidence:
    """Summarize native text evidence across the presentation."""
    return TextEvidence(
        total_characters=totals.total_characters,
        pages_or_units_with_text=len(totals.slide_numbers_with_text),
        estimated_text_density=totals.total_characters / max(total_slides, 1),
        native_text_available=totals.total_characters > 0,
    )


def build_table_evidence(totals: PptxFeatureTotals) -> TableEvidence:
    """Summarize table evidence across the presentation."""
    return TableEvidence(
        count=totals.number_of_tables,
        pages_or_units_with_tables=len(totals.slide_numbers_with_tables),
        # PPTX stores no semantic size metadata for table shapes beyond their EMU
        # dimensions.  A large/small split is not reliable at this extraction stage.
        large_count=0,
    )


def build_visual_evidence(totals: PptxFeatureTotals) -> VisualEvidence:
    """Summarize image, chart, drawing, and caption evidence across the presentation."""
    return VisualEvidence(
        embedded_image_count=totals.number_of_images,
        large_embedded_image_count=totals.number_of_large_images,
        # Drawing shapes (rectangles, lines, arrows, groups) map to vector_graphics_count
        # because they are vector geometry — not raster images.
        vector_graphics_count=totals.number_of_drawing_shapes,
        chart_count=totals.number_of_charts,
        # PPTX does not embed SVGs; vector art is stored as DrawingML (counted above).
        svg_count=0,
        pages_or_units_with_visuals=len(totals.slide_numbers_with_visuals),
        captioned_visual_count=totals.number_of_captioned_figures,
    )


def choose_visual_candidates_for_routing(
    visual_candidates: list[VisualCandidate],
    *,
    max_candidates_to_keep: int,
) -> list[VisualCandidate]:
    """
    Keep the visual elements most useful for engine routing.

    A dense presentation may produce many candidates (one per image shape, chart,
    diagram-heavy slide cluster).  The router only needs the strongest examples.
    """
    ranked = sorted(visual_candidates, key=_candidate_routing_priority, reverse=True)
    return ranked[:max_candidates_to_keep]


def _candidate_routing_priority(candidate: VisualCandidate) -> tuple[int, float]:
    """Score a visual candidate for routing relevance."""
    caption_bonus = 1 if (candidate.nearby_text or candidate.caption_or_alt_text) else 0
    return (caption_bonus + len(candidate.evidence), candidate.area_ratio or 0.0)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def calculate_shape_area_as_fraction_of_slide(
    shape: Any, slide_area_in_emus: float
) -> float | None:
    """
    Return how much of the slide this shape covers, expressed as a fraction
    between 0.0 and 1.0.

    Shape dimensions are stored in EMUs.  We don't need to convert to inches or
    pixels — we only need the ratio, so dividing shape area by slide area gives
    a unit-independent fraction.

    Returns None if the shape dimensions cannot be read.
    """
    try:
        width = float(shape.width or 0)
        height = float(shape.height or 0)
    except Exception:
        return None
    return max(width * height, 0.0) / max(slide_area_in_emus, 1.0)


def shape_is_a_drawing_or_connector(shape: Any, mso_shape_type: Any) -> bool:
    """
    Return True if this shape is a drawn element — a rectangle, line, arrow,
    freeform polygon, or grouped drawing — rather than a text box or table.

    Text boxes and tables are technically "shapes" in PPTX but carry structured
    content handled separately (Steps 2a and 2b).  We only want shapes that are
    purely visual — drawn geometry with no structured content of their own.

    Shape type constants used:
        AUTO_SHAPE  — built-in shapes like rectangles, circles, arrows, stars
        FREEFORM    — hand-drawn or custom polygon paths
        GROUP       — multiple shapes grouped together (often used for diagrams)
        LINE        — straight connector lines and arrows
    """
    if getattr(shape, "has_text_frame", False) or getattr(shape, "has_table", False):
        return False
    return shape.shape_type in {
        mso_shape_type.AUTO_SHAPE,
        mso_shape_type.FREEFORM,
        mso_shape_type.GROUP,
        mso_shape_type.LINE,
    }


def find_caption_in_slide_text(text_blocks_so_far: list[str]) -> str | None:
    """
    Search text collected so far on the slide for a figure caption line.

    If no caption is found, fall back to the first text block (usually the slide
    title) as loose contextual information for a downstream vision model.

    "So far" matters: this function is called mid-loop, so it only sees text
    from shapes processed before the current shape.  That is intentional —
    captions usually appear near or below the visual they describe.
    """
    for text in text_blocks_so_far:
        if contains_figure_caption(text):
            return compact_text(text)
    return compact_text(text_blocks_so_far[0]) if text_blocks_so_far else None


# ---------------------------------------------------------------------------
# File opening
# ---------------------------------------------------------------------------


def open_presentation(path: Path) -> tuple[Any, Any]:
    """
    Import python-pptx, open the file at *path*, and return the presentation
    object together with the MSO_SHAPE_TYPE enum.

    Both are returned together so the caller does not need to import python-pptx
    directly — the deferred import stays in one place.

    Raises
    ------
    DocumentError
        If python-pptx is not installed or the file cannot be opened.
    """
    try:
        from pptx import Presentation  # noqa: PLC0415 — intentional deferred import
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: PLC0415
    except ImportError as exc:
        raise DocumentError(
            "python-pptx is required for PPTX feature extraction. "
            "Install it with: pip install python-pptx",
            context={"path": str(path)},
        ) from exc

    try:
        return Presentation(str(path)), MSO_SHAPE_TYPE
    except Exception as exc:
        raise DocumentError(
            f"Could not open PowerPoint file for feature extraction: {path.name}",
            context={"path": str(path)},
        ) from exc
